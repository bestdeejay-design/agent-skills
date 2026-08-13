#!/usr/bin/env python3
"""Build a real .pptx deck from a JSON spec produced by the presentation-maker skill.

Usage:
    python build_pptx.py deck.json out.pptx
    python build_pptx.py --demo demo.pptx        # smoke-test deck

JSON spec shape:
{
  "title": "Deck title",
  "theme": {"name": "general", "palette": {"primary": "#007AFF", "background": "#FFFFFF",
             "card": "#F5F5F7", "stroke": "#E5E5EA", "background_text": "#FFFFFF",
             "primary_text": "#1C1C1E", "graph_0": "#007AFF", "graph_1": "#30B0C7"}},
  "slides": [
    {"type": "title",     "title": "...", "subtitle": "...", "presenter": "...", "date": "..."},
    {"type": "bullets",   "title": "...", "bullets": ["...", "..."]},
    {"type": "comparison","title": "...", "columns": [{"heading": "A", "points": [...]}, ...]},
    {"type": "metrics",   "title": "...", "metrics": [{"value": "200+", "label": "..."} , ...]},
    {"type": "table",     "title": "...", "table": {"headers": [...], "rows": [[...]]}},
    {"type": "chart",     "title": "...", "chart": {"title": "...", "categories": [...],
                             "series": [{"name": "...", "values": [...]}]}},
    {"type": "process",   "title": "...", "steps": ["...", "..."]},
    {"type": "divider",   "title": "...", "subtitle": "..."},
    {"type": "closing",   "title": "...", "subtitle": "..."}
  ]
}

Run with any modern python3 + python-pptx (pip install python-pptx).
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def _hex_to_rgb(hex_str: str, fallback=(0x1C, 0x1C, 0x1E)):
    """#RRGGBB -> (r,g,b) ints; tolerant of missing #."""
    s = hex_str.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    try:
        return tuple(int(s[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return fallback


class Palette:
    def __init__(self, palette: dict):
        self.primary = _hex_to_rgb(palette.get("primary", "#007AFF"))
        self.background = _hex_to_rgb(palette.get("background", "#FFFFFF"), (0xFF, 0xFF, 0xFF))
        self.card = _hex_to_rgb(palette.get("card", "#F5F5F7"), (0xF5, 0xF5, 0xF7))
        self.stroke = _hex_to_rgb(palette.get("stroke", "#E5E5EA"))
        self.accent_soft = _hex_to_rgb(palette.get("accent_soft", "#EAF0FF"))
        self.background_text = _hex_to_rgb(palette.get("background_text", "#FFFFFF"), (0xFF, 0xFF, 0xFF))
        self.primary_text = _hex_to_rgb(palette.get("primary_text", "#1C1C1E"), (0x1C, 0x1C, 0x1E))
        self.graphs = [
            _hex_to_rgb(palette.get(f"graph_{i}", "#007AFF"))
            for i in range(10)
        ]


def _set_slide_bg(slide, color):
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def _add_text(slide, left, top, width, height, text, size=20, bold=False,
              color=(0x1C, 0x1C, 0x1E), font="Inter", align="left", valign="top"):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    color = RGBColor(*color)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[valign]
    para = tf.paragraphs[0]
    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                      "right": PP_ALIGN.RIGHT}[align]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _add_rounded_card(slide, x, y, w, h, fill_rgb, stroke_rgb=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if stroke_rgb:
        shape.line.color.rgb = RGBColor(*stroke_rgb)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # soften corner radius
    try:
        shape.adjustments[0] = 0.06
    except Exception:
        pass
    return shape


def _render_title(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    accent = p.primary
    # top-left small logo bar + slide marker is skipped (kept minimal)
    _add_text(slide, 0, 0, W, H, spec.get("title", ""), size=44, bold=True,
              color=text_color, align="center", valign="middle")
    sub = spec.get("subtitle")
    if sub:
        _add_text(slide, 0, H // 2 + 30, W, 80, sub, size=22, color=accent,
                  align="center", valign="top")
    meta = "   ·   ".join(x for x in [spec.get("presenter"), spec.get("date")] if x)
    if meta:
        _add_text(slide, 0, H - 140, W, 40, meta, size=16, color=text_color,
                  align="center", valign="top")


def _render_divider(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.primary)
    _add_text(slide, 0, 0, W, H, spec.get("title", ""), size=44, bold=True,
              color=p.background_text, align="center", valign="middle")
    sub = spec.get("subtitle")
    if sub:
        _add_text(slide, 0, H // 2 + 30, W, 80, sub, size=22,
                  color=p.background_text, align="center", valign="top")


def _render_bullets(slide, spec, p: Palette, W, H, dark_bg: bool):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    bullets = spec.get("bullets", [])
    box = slide.shapes.add_textbox(100, 160, W - 200, H - 220)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(16)
        run = para.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(22)
        run.font.name = "Inter"
        run.font.color.rgb = RGBColor(*text_color)


def _render_comparison(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    columns = spec.get("columns", [])
    n = max(len(columns), 1)
    col_w = (W - 240) // n
    for idx, col in enumerate(columns):
        x = 120 + idx * (col_w + 20)
        card = _add_rounded_card(slide, x, 160, col_w, H - 240, p.card, p.stroke)
        _add_text(slide, x + 20, 180, col_w - 40, 50, col.get("heading", ""), size=24,
                  bold=True, color=p.primary, align="left")
        body = slide.shapes.add_textbox(x + 20, 240, col_w - 40, H - 400)
        tf = body.text_frame
        tf.word_wrap = True
        for i, point in enumerate(col.get("points", [])):
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = para.add_run()
            run.text = f"• {point}"
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(*text_color)


def _render_metrics(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    metrics = spec.get("metrics", [])
    n = max(len(metrics), 1)
    card_w = (W - 240 - (n - 1) * 20) // n
    for idx, m in enumerate(metrics):
        x = 120 + idx * (card_w + 20)
        _add_rounded_card(slide, x, 180, card_w, H - 300, p.card, p.stroke)
        _add_text(slide, x, 220, card_w, 120, str(m.get("value", "")), size=48, bold=True,
                  color=p.primary, align="center", valign="middle")
        _add_text(slide, x, 350, card_w, 80, str(m.get("label", "")), size=20,
                  color=text_color, align="center", valign="top")


def _render_table(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    t = spec.get("table", {})
    headers = t.get("headers", [])
    rows = t.get("rows", [])
    hl = t.get("highlight_col")
    n_rows = len(rows) + 1
    n_cols = max(len(headers), *(len(r) for r in rows), 1)
    from pptx.util import Inches, Pt as _Pt
    from pptx.dml.color import RGBColor
    table = slide.shapes.add_table(n_rows, n_cols, Inches(1), Inches(2.2),
                                   Inches(W - 2), Inches(H - 3.4)).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*p.primary)
        for par in cell.text_frame.paragraphs:
            for run in par.runs:
                run.font.color.rgb = RGBColor(*p.background_text)
                run.font.bold = True
    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            cell = table.cell(r, c)
            cell.text = str(val)
            if c == hl:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*p.accent_soft)
                for par in cell.text_frame.paragraphs:
                    for run in par.runs:
                        run.font.color.rgb = RGBColor(*p.primary)
                        run.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*p.card)


def _render_chart(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    ch = spec.get("chart", {})
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    chart_data = CategoryChartData()
    chart_data.categories = ch.get("categories", [])
    for i, s in enumerate(ch.get("series", [])):
        color = p.graphs[i % len(p.graphs)]
        chart_data.add_series(s.get("name", ""), s.get("values", []))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(2.2), Inches(W - 2), Inches(H - 3.4), chart_data,
    ).chart
    chart.has_legend = len(ch.get("series", [])) > 1
    # color series
    plot = chart.plots[0]
    try:
        from pptx.dml.color import RGBColor
        for idx, series in enumerate(plot.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*p.graphs[idx % len(p.graphs)])
    except Exception:
        pass
    try:
        from pptx.util import Pt
        chart.font.size = Pt(12)
    except Exception:
        pass


def _render_process(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    steps = spec.get("steps", [])
    n = max(len(steps), 1)
    step_w = (W - 200 - (n - 1) * 30) // n
    y = 220
    for i, step in enumerate(steps):
        x = 100 + i * (step_w + 30)
        _add_rounded_card(slide, x, y, step_w, H - 340, p.card, p.stroke)
        _add_text(slide, x, y + 16, step_w, 40, f"0{i + 1}", size=18, bold=True,
                  color=p.primary, align="center")
        _add_text(slide, x + 12, y + 70, step_w - 24, H - 440, step, size=18,
                  color=text_color, align="center", valign="top")
        if i < n - 1:
            _add_text(slide, x + step_w - 6, y + 30, 40, 40, "→", size=24,
                      color=p.primary, align="center")


def _render_closing(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 0, 0, W, H, spec.get("title", "Спасибо!"), size=48, bold=True,
              color=text_color, align="center", valign="middle")
    sub = spec.get("subtitle")
    if sub:
        _add_text(slide, 0, H // 2 + 40, W, 80, sub, size=22, color=p.primary,
                  align="center", valign="top")


def _render_feature(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _add_text(slide, 80, 60, W - 160, 80, spec.get("title", ""), size=34, bold=True, color=text_color)
    features = spec.get("features", [])
    n = max(len(features), 1)
    n_cols = 2 if n > 2 else n
    n_rows = (n + n_cols - 1) // n_cols
    gap = 24
    margin = 120
    usable_w = W - 2 * margin
    usable_h = H - 220
    card_w = (usable_w - (n_cols - 1) * gap) // n_cols
    card_h = (usable_h - (n_rows - 1) * gap) // n_rows
    for idx, f in enumerate(features):
        r, c = divmod(idx, n_cols)
        x = margin + c * (card_w + gap)
        y = 200 + r * (card_h + gap)
        _add_rounded_card(slide, x, y, card_w, card_h, p.card, p.stroke)
        title = f.get("title", "")
        if title:
            _add_text(slide, x + 24, y + 20, card_w - 48, 60, title, size=24,
                      bold=True, color=p.primary, align="left")
        body = f.get("text", "")
        if body:
            _add_text(slide, x + 24, y + 92, card_w - 48, card_h - 110, body, size=18,
                      color=text_color, align="left", valign="top")


def _render_big_number(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    accent = p.primary
    _add_text(slide, 0, H // 2 - 150, W, 160, str(spec.get("value", "")), size=110, bold=True,
              color=accent, align="center", valign="middle")
    label = spec.get("label", "")
    if label:
        _add_text(slide, 0, H // 2 + 40, W, 70, label, size=30, bold=True,
                  color=text_color, align="center", valign="top")
    sub = spec.get("subtitle", "")
    if sub:
        _add_text(slide, 0, H // 2 + 130, W, 110, sub, size=18,
                  color=text_color, align="center", valign="top")


RENDERERS = {
    "title": _render_title,
    "divider": _render_divider,
    "bullets": _render_bullets,
    "comparison": _render_comparison,
    "metrics": _render_metrics,
    "table": _render_table,
    "chart": _render_chart,
    "process": _render_process,
    "feature": _render_feature,
    "big_number": _render_big_number,
    "closing": _render_closing,
}


def build(spec: dict, out_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    W = Emu(prs.slide_width)
    H = Emu(prs.slide_height)
    palette = Palette(spec.get("theme", {}).get("palette", {}))
    dark_bg = _luma(palette.background) < 128

    for slide_spec in spec.get("slides", []):
        slide = prs.slides.add_slide(blank)
        stype = slide_spec.get("type", "bullets")
        renderer = RENDERERS.get(stype, _render_bullets)
        renderer(slide, slide_spec, palette, W, H, dark_bg)

    prs.save(str(out_path))
    return out_path


def _luma(rgb) -> float:
    return 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]


def _demo_spec() -> dict:
    return {
        "title": "UniverID — цифровая экосистема",
        "theme": {
            "name": "swift",
            "palette": {
                "primary": "#007AFF", "background": "#FFFFFF", "card": "#F5F5F7",
                "stroke": "#E5E5EA", "background_text": "#FFFFFF", "primary_text": "#1C1C1E",
                "graph_0": "#007AFF", "graph_1": "#30B0C7",
            },
        },
        "slides": [
            {"type": "title", "title": "UniverID", "subtitle": "Цифровая экосистема университета",
             "presenter": "AXIIOM", "date": "2026"},
            {"type": "divider", "title": "Проблема", "subtitle": "Цифровой хаос"},
            {"type": "bullets", "title": "Проблемы",
             "bullets": ["Разрозненные системы: Moodle, 1С, деканат",
                         "Нет единого окна для студента",
                         "Слабая коммуникация, потерянные письма",
                         "Нет персонализации под роль"]},
            {"type": "metrics", "title": "Масштаб",
             "metrics": [{"value": "20", "label": "модулей"},
                         {"value": "200+", "label": "возможностей"},
                         {"value": "6", "label": "ролей"},
                         {"value": "14", "label": "AI-функций"}]},
            {"type": "comparison", "title": "Сейчас → Цель",
             "columns": [{"heading": "Сейчас", "points": ["Разрозненные сервисы", "Нет единого входа", "Общий интерфейс"]},
                         {"heading": "UniverID", "points": ["Одно окно", "SSO + Госуслуги", "Ролевой интерфейс"]}]},
            {"type": "chart", "title": "Рост цифровизации",
             "chart": {"title": "Модулей по фазам",
                       "categories": ["Фаза 1", "Фаза 2", "Фаза 3"],
                       "series": [{"name": "Модули", "values": [8, 14, 20]}]}},
            {"type": "process", "title": "Дорожная карта",
             "steps": ["Фаза 1: Демо", "Фаза 2: Развитие", "Фаза 3: Экосистема"]},
            {"type": "closing", "title": "Давайте обсудим", "subtitle": "univerid.ru · AXIIOM"},
        ],
    }


def main() -> int:
    ap = argparse.ArgumentParser(description="Build .pptx from presentation-maker JSON spec")
    ap.add_argument("spec", nargs="?", help="JSON spec file")
    ap.add_argument("out", nargs="?", help="output .pptx path")
    ap.add_argument("--demo", help="write a demo deck to this path and exit")
    args = ap.parse_args()

    if args.demo:
        build(_demo_spec(), Path(args.demo))
        print(f"Demo deck written: {args.demo}")
        return 0

    if not args.spec or not args.out:
        ap.print_help()
        return 2

    try:
        spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    except Exception as e:
        print(f"Failed to read spec {args.spec}: {e}", file=sys.stderr)
        return 1

    out = build(spec, Path(args.out))
    print(f"Deck written: {out} ({len(spec.get('slides', []))} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
