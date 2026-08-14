#!/usr/bin/env python3
"""Build a real .pptx deck from a JSON spec produced by the presentation-maker skill.

Full design-system build: every slide carries the complete element set
(eyebrow rubricator, ghost index number, chrome logo + page indicator,
card shadows, numbered OVALs, accent metric values). Canvas is a 1600x900
design-pixel grid converted to EMU via PX=7620 — exact, no rounding gap.

Usage:
    python build_pptx.py deck.json out.pptx
    python build_pptx.py --demo demo.pptx        # smoke-test deck (all 14 types)

JSON spec shape:
{
  "theme": {"name": "general", "palette": {"primary": "#007AFF", "background": "#FFFFFF",
             "card": "#F5F5F7", "stroke": "#E5E5EA", "background_text": "#FFFFFF",
             "primary_text": "#1C1C1E", "muted": "#8B93A7", "accent_soft": "#EAF0FF",
             "graph_0": "#007AFF", "graph_1": "#30B0C7", ...}},
  "slides": [
    {"type": "title",  "eyebrow": "...", "title": "...", "subtitle": "...",
     "presenter": "...", "date": "..."},
    {"type": "divider", "eyebrow": "...", "title": "...", "subtitle": "..."},
    {"type": "bullets", "eyebrow": "...", "title": "...", "bullets": ["...", "..."]},
    {"type": "comparison", "eyebrow": "...", "title": "...",
     "columns": [{"heading": "A", "points": [...]}, ...]},
    {"type": "metrics", "eyebrow": "...", "title": "...",
     "metrics": [{"value": "200+", "label": "..."}, ...]},
    {"type": "table", "eyebrow": "...", "title": "...",
     "table": {"headers": [...], "rows": [[...]], "highlight_col": 1}},
    {"type": "chart", "eyebrow": "...", "title": "...",
     "chart": {"title": "...", "categories": [...],
               "series": [{"name": "...", "values": [...]}]}},
    {"type": "process", "eyebrow": "...", "title": "...", "steps": ["...", "..."]},
    {"type": "timeline", "eyebrow": "...", "title": "...",
     "items": [{"title": "...", "desc": "..."}, ...]},
    {"type": "feature", "eyebrow": "...", "title": "...",
     "features": [{"title": "...", "text": "..."}, ...]},
    {"type": "big_number", "eyebrow": "...", "title": "...",
     "value": "2,4×", "label": "...", "subtitle": "..."},
    {"type": "quote", "eyebrow": "...", "quote": "...", "attribution": "..."},
    {"type": "table_of_contents", "eyebrow": "...", "title": "...",
     "items": [{"title": "...", "desc": "..."}, ...]},
    {"type": "closing", "eyebrow": "...", "title": "...", "subtitle": "...",
     "presenter": "...", "date": "..."}
  ]
}

Run with any modern python3 + python-pptx (pip install python-pptx).
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt


def _hex_to_rgb(hex_str: str, fallback=(0x1C, 0x1C, 0x1E)):
    """#RRGGBB -> (r,g,b) ints; tolerant of missing #."""
    s = str(hex_str).strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    try:
        return tuple(int(s[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return fallback


# Design-space coordinate scale. slides.html renders at 1600x900 css px;
# the PPTX canvas is 13.333x7.5in = 12192000x6858000 EMU, so 1 px = 7620 EMU.
PX = 7620


def _px(n) -> "Emu":
    return Emu(int(n) * PX)


def _px_to_pt(px_val: float) -> float:
    """Design px -> points (1 pt = 1/72in, 1 px = 7620 EMU = 1/120 in)."""
    return px_val * 72.0 / 120.0


# ---- Layout contracts (design pixels) --------------------------------------
MARGIN_X = 96
TITLE_Y = 56
TITLE_H = 96          # title autofits to max 2 lines inside this band
CONTENT_Y = 180       # all content blocks start below the title band
BOTTOM_STOP = 80      # nothing may end lower than this


class Palette:
    def __init__(self, palette: dict):
        self.primary = _hex_to_rgb(palette.get("primary", "#007AFF"))
        self.background = _hex_to_rgb(palette.get("background", "#FFFFFF"), (0xFF, 0xFF, 0xFF))
        self.card = _hex_to_rgb(palette.get("card", "#F5F5F7"), (0xF5, 0xF5, 0xF7))
        self.stroke = _hex_to_rgb(palette.get("stroke", "#E5E5EA"))
        self.accent_soft = _hex_to_rgb(palette.get("accent_soft", "#EAF0FF"))
        self.background_text = _hex_to_rgb(palette.get("background_text", "#FFFFFF"), (0xFF, 0xFF, 0xFF))
        self.primary_text = _hex_to_rgb(palette.get("primary_text", "#1C1C1E"), (0x1C, 0x1C, 0x1E))
        self.muted = _hex_to_rgb(palette.get("muted", "#8B93A7"), (0x8B, 0x93, 0xA7))
        self.graphs = [
            _hex_to_rgb(palette.get(f"graph_{i}", "#007AFF"))
            for i in range(10)
        ]


def _set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def _luma(rgb) -> float:
    return 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]


# ---- Text fitting engine ---------------------------------------------------
# Arial glyph-width heuristics (em fractions). Good to ~10-15% for mixed
# Latin/Cyrillic text; the <a:normAutofit/> guard absorbs the residual error.
_WIDTH_NARROW = frozenset("iljftI.,:;!|()[]'\u2013\u2014-")
_WIDTH_WIDE = frozenset("wmWM\u0449\u0429\u0428\u0428\u042B\u042E\u042E")


def _char_w(ch: str) -> float:
    if ch == " ":
        return 0.30
    if ch in _WIDTH_NARROW:
        return 0.32
    if ch in _WIDTH_WIDE:
        return 0.88
    return 0.55


def _text_w_pt(text: str, size_pt: float) -> float:
    """Estimated rendered width in points (Arial)."""
    scale = size_pt / 72.0
    return sum(_char_w(c) for c in text) * 72.0 * scale


def _wrap_lines(text: str, size_pt: float, max_w_pt: float) -> list:
    """Greedy word wrap of a single paragraph onto max_w_pt width."""
    words = text.split()
    if not words:
        return [""]
    lines, cur, cur_w = [], "", 0.0
    space_w = _text_w_pt(" ", size_pt)
    for w in words:
        w_w = _text_w_pt(w, size_pt)
        add = (cur_w + space_w + w_w) if cur else w_w
        if cur and add > max_w_pt:
            lines.append(cur)
            cur, cur_w = w, w_w
        else:
            cur = (cur + " " + w) if cur else w
            cur_w = add
    if cur:
        lines.append(cur)
    return lines or [""]


def _fit_size(texts, max_w_pt, max_h_pt, start_size=20, min_size=10,
              line_h=1.2):
    """Binary-search the largest font size fitting all texts inside the box.

    Heuristic only; <a:normAutofit/> in the XML shrinks further in rendering.
    """
    lo, hi = min_size, start_size
    best = start_size
    while lo <= hi:
        mid = (lo + hi) // 2
        total_h = 0.0
        ok = True
        for t in texts:
            lines = _wrap_lines(t, mid, max_w_pt)
            total_h += len(lines) * mid * line_h
            if total_h > max_h_pt or any(
                _text_w_pt(ln, mid) > max_w_pt + 0.5 for ln in lines
            ):
                ok = False
                break
        if ok:
            best = mid
            lo = mid + 1
        else:
            hi = mid - 1
    return best


def _add_text(slide, left, top, width, height, text, size=20, bold=False,
              color=(0x1C, 0x1C, 0x1E), align="left", valign="top",
              min_size=10, line_h=1.2, alpha=None, tracking=None,
              fit=True, font="Arial", name=None):
    """Text box in design pixels with autofit; returns the TextBox.

    `alpha` — transparency percent (0 opaque .. 100 invisible), via <a:alpha>.
    `tracking` — letter spacing in points, via <a:spc> on the run.
    `name` — optional shape name (e.g. "quote-mark") so extern linters
             (qa_intern.py) can whitelist design elements by prefix.
    """
    box = slide.shapes.add_textbox(_px(left), _px(top), _px(width), _px(height))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[valign]
    para = tf.paragraphs[0]
    para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                      "right": PP_ALIGN.RIGHT}[align]
    run = para.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = font
    if alpha is not None:
        rPr = run._r.get_or_add_rPr()
        solid = rPr.makeelement(qn('a:solidFill'), {})
        srgb = rPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % color})
        a = rPr.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))})
        srgb.append(a)
        solid.append(srgb)
        rPr.insert(0, solid)
    if tracking is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(tracking * 100)))
        # Arial fallback keeps the same widths across Office/LibreOffice.
    if fit:
        max_w_pt = _px_to_pt(width) * 0.96
        max_h_pt = _px_to_pt(height)
        fitted = _fit_size([str(text)], max_w_pt, max_h_pt,
                           start_size=size, min_size=min_size, line_h=line_h)
        if fitted < size:
            run.font.size = Pt(fitted)
            size = fitted
        # embed normAutofit so PowerPoint/LibreOffice shrink the text if the
        # heuristic overestimated the width.
        try:
            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            norm = bodyPr.makeelement(qn('a:normAutofit'), {})
            bodyPr.append(norm)
        except Exception:
            pass
    return box


def _add_bullets(slide, left, top, width, height, items, color, size=22,
                 min_size=13, style="markers", accent=(0x00, 0x7A, 0xFF)):
    """Bullet list; style 'markers' = accent squares, 'cards' = white cards."""
    from pptx.util import Pt as _Pt
    if style == "cards":
        y = top
        item_h = 62
        gap = 14
        for i, b in enumerate(items):
            card_h = item_h
            box_h = item_h
            card = _add_rounded_card(slide, left, y, width, card_h,
                                     (0xFF, 0xFF, 0xFF), None)
            _add_shadow(card, blur_pt=8, dist_pt=2, alpha_pct=8)
            card.name = "bullets-card"
            mk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        _px(left + 18), _px(y + 24),
                                        _px(12), _px(12))
            mk.fill.solid()
            mk.fill.fore_color.rgb = RGBColor(*accent)
            mk.line.fill.background()
            mk.name = "bullet-marker"
            _add_text(slide, left + 46, y, width - 62, box_h, b, size=size,
                      color=color, align="left", valign="middle",
                      min_size=min_size, line_h=1.2)
            y += card_h + gap
        return
    # markers (default): accent squares in front of each item
    n = max(len(items), 1)
    per = (height - (n - 1) * 18) // n
    for i, b in enumerate(items):
        y = top + i * (per + 18)
        mk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    _px(left), _px(y + 12), _px(12), _px(12))
        mk.fill.solid()
        mk.fill.fore_color.rgb = RGBColor(*accent)
        mk.line.fill.background()
        mk.name = "bullet-marker"
        _add_text(slide, left + 30, y, width - 30, per, b, size=size,
                  color=color, align="left", valign="middle",
                  min_size=min_size, line_h=1.2)


def _add_rounded_card(slide, x, y, w, h, fill_rgb, stroke_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _px(x), _px(y), _px(w), _px(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if stroke_rgb:
        shape.line.color.rgb = RGBColor(*stroke_rgb)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.06
    except Exception:
        pass
    return shape


def _add_shadow(shape, blur_pt=8, dist_pt=2, alpha_pct=10):
    """Soft outer shadow (XML) — design-system '--shadow' equivalent."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    eff = spPr.makeelement(qn('a:effectLst'), {})
    shdw = eff.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur_pt * 12700)), 'dist': str(int(dist_pt * 12700)),
        'dir': '5400000', 'rotWithShape': '0'})
    clr = shdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    a = clr.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    clr.append(a)
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def _add_ghost_num(slide, num, color, W, H):
    """Editorial index number ('.index-num'): huge translucent digit behind content."""
    box = _add_text(slide, W - 420, -30, 380, 430, f"{num:02d}", size=300,
                    bold=True, color=color, align="right", valign="middle",
                    min_size=160, line_h=1.0, alpha=7)
    box.name = "ghost-num"
    return box


def _add_decor_oval(slide, x, y, w, h, fill_rgb, alpha_pct):
    """Semi-transparent decorative circle on dark slides."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, _px(x), _px(y), _px(w), _px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shp.line.fill.background()
    spPr = shp._element.spPr
    solid = spPr.find(qn('a:solidFill'))
    if solid is not None:
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            el = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
            srgb.append(el)
    shp.name = "decor"
    return shp


def _slide_chrome(slide, p: Palette, W, idx, total, dark: bool):
    """Logo badge + page indicator ('.slide-head' equivalent) on every slide."""
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  _px(48), _px(28), _px(30), _px(30))
    try:
        logo.adjustments[0] = 0.25
    except Exception:
        pass
    logo.fill.solid()
    logo.fill.fore_color.rgb = RGBColor(*(p.background_text if dark else p.primary))
    logo.line.fill.background()
    logo.name = "chrome-logo"
    tfo = logo.text_frame
    tfo.word_wrap = False
    po = tfo.paragraphs[0]
    po.alignment = PP_ALIGN.CENTER
    r0 = po.add_run()
    r0.text = "\u041d"
    r0.font.size = Pt(13)
    r0.font.bold = True
    r0.font.name = "Arial"
    r0.font.color.rgb = RGBColor(*(p.primary if dark else p.background_text))
    box = _add_text(slide, W - 250, 32, 202, 24, f"{idx:02d} / {total}", size=11,
                    bold=True, color=(0xC9, 0xDB, 0xDE) if dark else p.muted,
                    align="right", fit=False, tracking=1.4)
    box.name = "chrome-page"


def _content_title(slide, spec, text_color, accent):
    eb = spec.get("eyebrow")
    if eb:
        _add_text(slide, MARGIN_X, 26, 1600 - 2 * MARGIN_X - 320, 24, str(eb), size=11,
                  bold=True, color=accent, align="left", fit=False, tracking=3.2)
    _add_text(slide, MARGIN_X, TITLE_Y, 1600 - 2 * MARGIN_X, TITLE_H,
              spec.get("title", ""), size=34, bold=True, color=text_color,
              min_size=18)


def _set_notes(slide, slide_spec):
    notes = slide_spec.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


# ---- Renderers -------------------------------------------------------------
def _render_title(slide, spec, p: Palette, W, H, dark_bg: bool):
    # Dark cover: primary background, light text (design-system "бутерброд").
    _set_slide_bg(slide, p.primary)
    white = p.background_text
    _add_decor_oval(slide, W - 400, -170, 500, 500, p.graphs[1], 22)
    _add_decor_oval(slide, -90, H - 240, 240, 240, p.graphs[1], 16)
    eb = spec.get("eyebrow")
    if eb:
        _add_text(slide, MARGIN_X, 88, W - 2 * MARGIN_X, 22, str(eb), size=12,
                  bold=True, color=p.accent_soft, align="center", fit=False,
                  tracking=3.6)
    _add_text(slide, MARGIN_X, 130, W - 2 * MARGIN_X, 330, spec.get("title", ""),
              size=50, bold=True, color=white, align="center", valign="middle",
              min_size=24, line_h=1.12)
    sub = spec.get("subtitle")
    if sub:
        _add_text(slide, MARGIN_X, 500, W - 2 * MARGIN_X, 90, sub,
                  size=20, color=white, align="center", min_size=14, alpha=82)
    meta = " \u00b7 ".join(x for x in [spec.get("presenter"), spec.get("date")] if x)
    if meta:
        _add_text(slide, MARGIN_X, H - 120, W - 2 * MARGIN_X, 40, meta, size=15,
                  color=white, align="center", min_size=12, alpha=70)


def _render_divider(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.primary)
    white = p.background_text
    _add_decor_oval(slide, W - 380, -140, 460, 460, p.graphs[1], 20)
    _add_decor_oval(slide, -80, H - 220, 220, 220, p.graphs[1], 15)
    eb = spec.get("eyebrow")
    if eb:
        _add_text(slide, MARGIN_X, 88, W - 2 * MARGIN_X, 22, str(eb), size=12,
                  bold=True, color=p.accent_soft, align="center", fit=False,
                  tracking=3.6)
    _add_text(slide, MARGIN_X, 130, W - 2 * MARGIN_X, 330, spec.get("title", ""),
              size=50, bold=True, color=white, align="center", valign="middle",
              min_size=24, line_h=1.12)
    sub = spec.get("subtitle")
    if sub:
        _add_text(slide, MARGIN_X, 500, W - 2 * MARGIN_X, 90, sub,
                  size=20, color=white, align="center", min_size=14, alpha=82)


def _render_bullets(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    _add_bullets(slide, 120, CONTENT_Y, W - 240, H - CONTENT_Y - BOTTOM_STOP,
                 spec.get("bullets", []), text_color, size=22, min_size=13,
                 style="cards", accent=p.primary)


def _render_comparison(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    columns = spec.get("columns", [])
    n = max(len(columns), 1)
    col_w = (W - 240 - (n - 1) * 20) // n
    for idx, col in enumerate(columns):
        x = 120 + idx * (col_w + 20)
        card = _add_rounded_card(slide, x, CONTENT_Y, col_w, H - CONTENT_Y - BOTTOM_STOP,
                                 p.card, p.stroke)
        _add_shadow(card, blur_pt=10, dist_pt=3, alpha_pct=9)
        card.name = "comparison-card"
        _add_text(slide, x + 24, CONTENT_Y + 20, col_w - 48, 52,
                  col.get("heading", ""), size=22, bold=True, color=p.primary,
                  min_size=15)
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    _px(x + 24), _px(CONTENT_Y + 82),
                                    _px(col_w - 48), _px(2))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(*p.stroke)
        ln.line.fill.background()
        ln.name = "hairline"
        _add_bullets(slide, x + 24, CONTENT_Y + 100, col_w - 48,
                     H - CONTENT_Y - BOTTOM_STOP - 120,
                     col.get("points", []), text_color, size=19, min_size=11,
                     accent=p.primary)


def _render_metrics(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    metrics = spec.get("metrics", [])
    n = max(len(metrics), 1)
    card_w = (W - 240 - (n - 1) * 20) // n
    for idx, m in enumerate(metrics):
        x = 120 + idx * (card_w + 20)
        card = _add_rounded_card(slide, x, CONTENT_Y, card_w, H - CONTENT_Y - BOTTOM_STOP,
                                 p.card, p.stroke)
        _add_shadow(card, blur_pt=10, dist_pt=3, alpha_pct=9)
        card.name = "metric-card"
        accent = p.graphs[idx % 8]
        _add_text(slide, x, CONTENT_Y + 44, card_w, 150, str(m.get("value", "")),
                  size=64, bold=True, color=accent, align="center",
                  valign="middle", min_size=30)
        _add_text(slide, x + 24, CONTENT_Y + 220, card_w - 48, 110,
                  str(m.get("label", "")), size=16, color=p.muted,
                  align="center", min_size=12, line_h=1.2)


def _render_table(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    t = spec.get("table", {})
    headers = [str(h) for h in t.get("headers", [])]
    rows = [[str(v) for v in r] for r in t.get("rows", [])]
    hl = t.get("highlight_col")
    n_rows = len(rows) + 1
    n_cols = max(len(headers), *(len(r) for r in rows), 1)
    tbl_w = W - 2 * MARGIN_X
    max_h = H - CONTENT_Y - BOTTOM_STOP
    first_w = int(tbl_w * 0.40)
    rest_w = max((tbl_w - first_w) // max(n_cols - 1, 1), 180)
    col_widths = [first_w if c == 0 else rest_w for c in range(n_cols)]

    def _cell_fit_size(texts, col_w_px, start, min_size=9):
        return _fit_size(
            texts,
            max_w_pt=_px_to_pt(col_w_px - 32) * 0.94,
            max_h_pt=_px_to_pt(420),
            start_size=start,
            min_size=min_size,
            line_h=1.2,
        )

    head_size = min(
        _cell_fit_size([headers[c]], col_widths[c], 15)
        for c in range(min(len(headers), n_cols))
    ) if headers else 15
    body_sizes = []
    for c in range(n_cols):
        texts = [r[c] for r in rows if c < len(r)]
        body_sizes.append(_cell_fit_size(texts, col_widths[c], 14) if texts else 14)

    def _est_h(b_size):
        total = 0.0
        for r_i, row in enumerate(rows):
            line_max = 1
            for c in range(n_cols):
                txt = row[c] if c < len(row) else ""
                lines = len(_wrap_lines(txt, b_size, _px_to_pt(col_widths[c] - 32) * 0.94))
                line_max = max(line_max, lines)
            total += line_max * b_size * 1.35 * 1.6667
        total += 1 * head_size * 1.35 * 1.6667 + 8
        return total

    while body_sizes and min(body_sizes) > 9 and _est_h(min(body_sizes)) > max_h:
        body_sizes = [s - 1 for s in body_sizes]

    table = slide.shapes.add_table(n_rows, n_cols, _px(MARGIN_X), _px(CONTENT_Y),
                                   _px(tbl_w), _px(max_h)).table
    for c in range(n_cols):
        table.columns[c].width = Emu(_px(col_widths[c]))
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*p.primary)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Emu(_px(16))
        cell.margin_right = Emu(_px(16))
        cell.margin_top = Emu(_px(6))
        cell.margin_bottom = Emu(_px(6))
        for par in cell.text_frame.paragraphs:
            for run in par.runs:
                run.font.name = "Arial"
                run.font.size = Pt(head_size)
                run.font.color.rgb = RGBColor(*p.background_text)
                run.font.bold = True
        table.rows[0].height = Emu(_px(int(1.35 * head_size * 1.6667) + 12))
    for r, row in enumerate(rows, start=1):
        est_lines = 1
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            cell = table.cell(r, c)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Emu(_px(16))
            cell.margin_right = Emu(_px(16))
            cell.margin_top = Emu(_px(6))
            cell.margin_bottom = Emu(_px(6))
            size = body_sizes[c] if c < len(body_sizes) else 14
            for par in cell.text_frame.paragraphs:
                for run in par.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(size)
                    run.font.color.rgb = RGBColor(*text_color)
            if c == hl:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*p.accent_soft)
                for par in cell.text_frame.paragraphs:
                    for run in par.runs:
                        run.font.color.rgb = RGBColor(*p.primary)
                        run.font.bold = True
            elif c == 0:
                for par in cell.text_frame.paragraphs:
                    for run in par.runs:
                        run.font.color.rgb = RGBColor(*p.primary)
                        run.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*p.card)
            n_lines = len(_wrap_lines(val, size, _px_to_pt(col_widths[c] - 32) * 0.94))
            est_lines = max(est_lines, n_lines)
        table.rows[r].height = Emu(_px(int(est_lines * 1.35 * size * 1.6667) + 12))


def _render_chart(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    ch = spec.get("chart", {})
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    series = ch.get("series", [])
    chart_data = CategoryChartData()
    chart_data.categories = ch.get("categories", [])
    for s in series:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        _px(MARGIN_X), _px(CONTENT_Y), _px(W - 2 * MARGIN_X),
        _px(H - CONTENT_Y - BOTTOM_STOP), chart_data).chart
    chart_title = ch.get("title")
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(chart_title)
        for par in chart.chart_title.text_frame.paragraphs:
            for run in par.runs:
                run.font.size = Pt(16)
                run.font.bold = True
    chart.has_legend = len(series) > 1
    if len(series) > 1:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.gap_width = 80
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(11)
    all_int = all(
        isinstance(v, int)
        for s in series for v in s.get("values", [])
    )
    dl.number_format = "0" if all_int else "0.0"
    dl.number_format_is_linked = False
    for idx, s in enumerate(plot.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = RGBColor(*p.graphs[idx % len(p.graphs)])
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.font.size = Pt(12)


def _render_process(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    steps = spec.get("steps", [])
    n = max(len(steps), 1)
    step_w = (W - 200 - (n - 1) * 30) // n
    y = CONTENT_Y + 40
    for i, step in enumerate(steps):
        x = 100 + i * (step_w + 30)
        card = _add_rounded_card(slide, x, y, step_w, H - y - BOTTOM_STOP, p.card, p.stroke)
        _add_shadow(card, blur_pt=10, dist_pt=3, alpha_pct=9)
        card.name = "process-card"
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   _px(x + step_w // 2 - 20), _px(y + 20),
                                   _px(40), _px(40))
        o.fill.solid()
        o.fill.fore_color.rgb = RGBColor(*p.primary)
        o.line.fill.background()
        o.name = "step-num"
        tfo = o.text_frame
        tfo.word_wrap = False
        po = tfo.paragraphs[0]
        po.alignment = PP_ALIGN.CENTER
        r0 = po.add_run()
        r0.text = f"{i + 1}"
        r0.font.size = Pt(17)
        r0.font.bold = True
        r0.font.name = "Arial"
        r0.font.color.rgb = RGBColor(*p.background_text)
        _add_text(slide, x + 16, y + 78, step_w - 32, H - y - BOTTOM_STOP - 100,
                  step, size=17, color=text_color, align="center", min_size=11,
                  line_h=1.25)
        if i < n - 1:
            _add_text(slide, x + step_w - 8, y + 30, 48, 44, "\u2192", size=22,
                      color=p.primary, align="center", min_size=16)


def _render_timeline(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    items = spec.get("items", [])
    n = max(len(items), 1)
    item_w = (W - 240 - (n - 1) * 30) // n
    y = CONTENT_Y + 40
    for i, it in enumerate(items):
        x = 120 + i * (item_w + 30)
        card = _add_rounded_card(slide, x, y, item_w, H - y - BOTTOM_STOP, p.card, p.stroke)
        _add_shadow(card, blur_pt=10, dist_pt=3, alpha_pct=9)
        card.name = "timeline-card"
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   _px(x + item_w // 2 - 18), _px(y + 18),
                                   _px(36), _px(36))
        o.fill.solid()
        o.fill.fore_color.rgb = RGBColor(*p.primary)
        o.line.fill.background()
        o.name = "step-num"
        tfo = o.text_frame
        tfo.word_wrap = False
        po = tfo.paragraphs[0]
        po.alignment = PP_ALIGN.CENTER
        r0 = po.add_run()
        r0.text = f"{i + 1}"
        r0.font.size = Pt(16)
        r0.font.bold = True
        r0.font.name = "Arial"
        r0.font.color.rgb = RGBColor(*p.background_text)
        _add_text(slide, x + 16, y + 72, item_w - 32, 52, it.get("title", ""),
                  size=18, bold=True, color=text_color, align="center", min_size=12)
        _add_text(slide, x + 16, y + 132, item_w - 32, H - y - BOTTOM_STOP - 160, it.get("desc", ""),
                  size=16, color=text_color, align="center", min_size=10)


def _render_closing(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.primary)
    white = p.background_text
    _add_decor_oval(slide, W - 340, -120, 420, 420, p.graphs[1], 22)
    _add_decor_oval(slide, -70, H - 200, 210, 210, p.graphs[1], 16)
    eb = spec.get("eyebrow")
    if eb:
        _add_text(slide, MARGIN_X, 88, W - 2 * MARGIN_X, 22, str(eb), size=12,
                  bold=True, color=p.accent_soft, align="center", fit=False,
                  tracking=3.6)
    _add_text(slide, MARGIN_X, 130, W - 2 * MARGIN_X, 330, spec.get("title", "\u0421\u043f\u0430\u0441\u0438\u0431\u043e!"),
              size=50, bold=True, color=white, align="center", valign="middle",
              min_size=24, line_h=1.12)
    sub = spec.get("subtitle")
    if sub:
        _add_text(slide, MARGIN_X, 500, W - 2 * MARGIN_X, 90, sub,
                  size=20, color=white, align="center", min_size=14, alpha=82)
    meta = " \u00b7 ".join(x for x in [spec.get("presenter"), spec.get("date")] if x)
    if meta:
        _add_text(slide, MARGIN_X, H - 120, W - 2 * MARGIN_X, 40, meta, size=15,
                  color=white, align="center", min_size=12, alpha=70)


def _render_feature(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    features = spec.get("features", [])
    n = max(len(features), 1)
    n_cols = 2 if n > 2 else n
    n_rows = (n + n_cols - 1) // n_cols
    gap = 24
    margin = 120
    usable_w = W - 2 * margin
    usable_h = H - CONTENT_Y - 40 - BOTTOM_STOP
    card_w = (usable_w - (n_cols - 1) * gap) // n_cols
    card_h = (usable_h - (n_rows - 1) * gap) // n_rows
    for idx, f in enumerate(features):
        r, c = divmod(idx, n_cols)
        x = margin + c * (card_w + gap)
        y = CONTENT_Y + 20 + r * (card_h + gap)
        card = _add_rounded_card(slide, x, y, card_w, card_h, p.card, p.stroke)
        _add_shadow(card, blur_pt=10, dist_pt=3, alpha_pct=9)
        card.name = "feature-card"
        title = f.get("title", "")
        if title:
            _add_text(slide, x + 24, y + 20, card_w - 48, 60, title, size=24,
                      bold=True, color=p.primary, min_size=15)
        body = f.get("text", "")
        if body:
            _add_text(slide, x + 24, y + 92, card_w - 48, card_h - 110, body,
                      size=18, color=text_color, min_size=11)


def _render_big_number(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    title = spec.get("title")
    if title:
        _content_title(slide, spec, text_color, p.primary)
    _add_text(slide, MARGIN_X, H // 2 - 150, W - 2 * MARGIN_X, 160,
              str(spec.get("value", "")), size=110, bold=True, color=p.primary,
              align="center", valign="middle", min_size=40)
    label = spec.get("label", "")
    if label:
        _add_text(slide, MARGIN_X, H // 2 + 40, W - 2 * MARGIN_X, 70, label,
                  size=30, bold=True, color=text_color, align="center", min_size=16)
    sub = spec.get("subtitle", "")
    if sub:
        _add_text(slide, MARGIN_X, H // 2 + 130, W - 2 * MARGIN_X, 110, sub,
                  size=18, color=text_color, align="center", min_size=12)


def _render_quote(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    eb = spec.get("eyebrow")
    if eb:
        _add_text(slide, MARGIN_X, 40, 1600 - 2 * MARGIN_X - 320, 22, str(eb), size=11,
                  bold=True, color=p.primary, align="left", fit=False,
                  tracking=3.2)
    _add_text(slide, 0, 130, W, 150, "\u201c", size=120, bold=True,
              color=p.primary, align="center", valign="top", min_size=80,
              name="quote-mark")
    quote = spec.get("quote", "")
    if quote:
        _add_text(slide, 160, 280, W - 320, 300, quote, size=34, bold=True,
                  color=text_color, align="center", valign="middle", min_size=17)
    attrib = spec.get("attribution") or spec.get("author") or ""
    if attrib:
        _add_text(slide, MARGIN_X, H - 160, W - 2 * MARGIN_X, 60, attrib, size=20,
                  color=p.muted, align="center", min_size=13)


def _render_table_of_contents(slide, spec, p: Palette, W, H, dark_bg: bool):
    _set_slide_bg(slide, p.background)
    text_color = p.background_text if dark_bg else p.primary_text
    _content_title(slide, spec, text_color, p.primary)
    items = spec.get("items", [])
    n = max(len(items), 1)
    n_cols = 2 if n > 4 else 1
    n_rows = (n + n_cols - 1) // n_cols
    gap = 40
    margin_x = 140
    margin_y = CONTENT_Y + 20
    usable_w = W - 2 * margin_x
    usable_h = H - margin_y - BOTTOM_STOP
    col_w = (usable_w - (n_cols - 1) * gap) // n_cols
    row_h = (usable_h - (n_rows - 1) * 48) // n_rows
    for idx, it in enumerate(items):
        r, c = divmod(idx, n_cols)
        x = margin_x + c * (col_w + gap)
        y = margin_y + r * (row_h + 48)
        o = slide.shapes.add_shape(MSO_SHAPE.OVAL, _px(x), _px(y), _px(44), _px(44))
        o.fill.solid()
        o.fill.fore_color.rgb = RGBColor(*p.primary)
        o.line.fill.background()
        o.name = "toc-num"
        tfo = o.text_frame
        tfo.word_wrap = False
        po = tfo.paragraphs[0]
        po.alignment = PP_ALIGN.CENTER
        r0 = po.add_run()
        r0.text = f"{idx + 1:02d}"
        r0.font.size = Pt(15)
        r0.font.bold = True
        r0.font.name = "Arial"
        r0.font.color.rgb = RGBColor(*p.background_text)
        _add_text(slide, x + 66, y + 2, col_w - 66, 40, it.get("title", ""), size=24,
                  bold=True, color=text_color, align="left", valign="middle",
                  min_size=14, name="toc-title")
        desc = it.get("desc", "")
        if desc:
            _add_text(slide, x + 66, y + 46, col_w - 66, 52, desc, size=17,
                      color=p.muted, align="left", min_size=12, name="toc-desc")


RENDERERS = {
    "title": _render_title,
    "divider": _render_divider,
    "bullets": _render_bullets,
    "comparison": _render_comparison,
    "metrics": _render_metrics,
    "table": _render_table,
    "chart": _render_chart,
    "process": _render_process,
    "timeline": _render_timeline,
    "feature": _render_feature,
    "big_number": _render_big_number,
    "quote": _render_quote,
    "table_of_contents": _render_table_of_contents,
    "closing": _render_closing,
}


def build(spec: dict, out_path: Path) -> Path:
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Emu(1600 * PX)
    prs.slide_height = Emu(900 * PX)
    blank = prs.slide_layouts[6]  # blank layout

    W, H = 1600, 900
    palette = Palette(spec.get("theme", {}).get("palette", {}))
    dark_bg = _luma(palette.background) < 128

    slides_spec = spec.get("slides", [])
    total = len(slides_spec)
    for idx, slide_spec in enumerate(slides_spec):
        slide = prs.slides.add_slide(blank)
        stype = slide_spec.get("type", "bullets")
        dark_slide = stype in ("title", "closing", "divider") or dark_bg
        _add_ghost_num(slide, idx + 1,
                       palette.background_text if dark_slide else palette.primary,
                       W, H)
        renderer = RENDERERS.get(stype, _render_bullets)
        renderer(slide, slide_spec, palette, W, H, dark_bg)
        _set_notes(slide, slide_spec)
        _slide_chrome(slide, palette, W, idx + 1, total, dark_slide)

    prs.save(str(out_path))
    return out_path


def _demo_spec() -> dict:
    return {
        "title": "Full element set demo (all 14 types)",
        "theme": {
            "palette": {
                "primary": "#0E4D5C", "background": "#F2F4F3", "card": "#FFFFFF",
                "stroke": "#D7DEDC", "background_text": "#FFFFFF",
                "primary_text": "#0E4D5C", "muted": "#5C6B70",
                "accent_soft": "#DCEAE9",
                "graph_0": "#0E4D5C", "graph_1": "#E98A5C", "graph_2": "#349886",
                "graph_3": "#9271B4", "graph_4": "#DBAD54", "graph_5": "#5482A0",
            },
        },
        "slides": [
            {"type": "title", "eyebrow": "NEO CONSULTING · АВГУСТ 2026",
             "title": "Стратегия цифровой трансформации",
             "subtitle": "Дорожная карта на 18 месяцев",
             "presenter": "Сергей Кузюков", "date": "14 августа 2026"},
            {"type": "table_of_contents", "eyebrow": "ПОВЕСТКА",
             "title": "Что обсудим сегодня",
             "items": [
                 {"title": "Контекст и цели", "desc": "Почему трансформация нужна сейчас"},
                 {"title": "Текущее состояние", "desc": "Аудит процессов и систем"},
                 {"title": "План внедрения", "desc": "Этапы, сроки, ответственные"},
                 {"title": "Риски и бюджет", "desc": "Что может пойти не так"}]},
            {"type": "divider", "eyebrow": "ЧАСТЬ 01",
             "title": "Контекст и цели",
             "subtitle": "Почему цифровая трансформация — вопрос выживания"},
            {"type": "bullets", "eyebrow": "КОНТЕКСТ",
             "title": "Три драйвера изменений",
             "bullets": [
                 "Операционные затраты выросли на 18% за два года",
                 "Конкуренты сокращают цикл вывода продукта вдвое",
                 "Клиенты ожидают цифровой опыт — 72% готовы уйти без него"]},
            {"type": "comparison", "eyebrow": "АНАЛИЗ",
             "title": "Сейчас и через 18 месяцев",
             "columns": [
                 {"heading": "Как сегодня",
                  "points": ["Ручной ввод данных, до 40% ошибок",
                             "Ответ на заявку — 3 дня",
                             "Отчётность готовится неделями"]},
                 {"heading": "Целевая модель",
                  "points": ["Автоматизация 80% документооборота",
                             "Ответ — в течение часа",
                             "Дашборды в реальном времени"]}]},
            {"type": "metrics", "eyebrow": "КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ",
             "title": "Целевые метрики программы",
             "metrics": [
                 {"value": "-38%", "label": "операционных затрат"},
                 {"value": "80%", "label": "документов без ручного ввода"},
                 {"value": "2,4×", "label": "скорость обработки заявок"},
                 {"value": "4,9", "label": "NPS клиентов (с 4,1)"}]},
            {"type": "table", "eyebrow": "АУДИТ",
             "title": "Зрелость процессов по направлениям",
             "table": {
                 "headers": ["Направление", "Текущий уровень", "Целевой уровень", "Разрыв"],
                 "rows": [
                     ["Документооборот", "1,5", "4,0", "2,5"],
                     ["Управление заявками", "2,0", "4,0", "2,0"],
                     ["Отчётность", "2,5", "4,5", "2,0"],
                     ["Клиентский сервис", "3,0", "4,5", "1,5"]],
                 "highlight_col": 2}},
            {"type": "chart", "eyebrow": "ФИНАНСЫ",
             "title": "Инвестиции по кварталам, млн ₽",
             "chart": {
                 "categories": ["Q3'26", "Q4'26", "Q1'27", "Q2'27", "Q3'27"],
                 "series": [
                     {"name": "Платформа", "values": [12, 18, 22, 18, 10]},
                     {"name": "Интеграции", "values": [6, 10, 14, 16, 8]}]}},
            {"type": "process", "eyebrow": "ДОРОЖНАЯ КАРТА",
             "title": "Четыре волны внедрения",
             "steps": [
                 "Аудит и проектирование — квартал 1",
                 "Пилот на одном подразделении — квартал 2",
                 "Масштабирование — кварталы 3–5",
                 "Оптимизация и поддержка — квартал 6"]},
            {"type": "timeline", "eyebrow": "ВЕХИ",
             "title": "Контрольные точки программы",
             "items": [
                 {"title": "Ноябрь 2026", "desc": "Пилот введён в эксплуатацию"},
                 {"title": "Март 2027", "desc": "80% процессов автоматизировано"},
                 {"title": "Август 2027", "desc": "Программа завершена"}]},
            {"type": "feature", "eyebrow": "РЕШЕНИЕ",
             "title": "Что даст платформа",
             "features": [
                 {"title": "Единое окно", "text": "Все процессы и документы в одном интерфейсе"},
                 {"title": "Умный ввод", "text": "Распознавание документов без ручного переноса"},
                 {"title": "Живая аналитика", "text": "Дашборды без участия аналитиков"},
                 {"title": "Интеграции", "text": "Готовая связка с 1С и банковскими сервисами"}]},
            {"type": "big_number", "eyebrow": "ОЖИДАЕМЫЙ ЭФФЕКТ",
             "title": "Сколько вернёт программа",
             "value": "2,4×", "label": "ROI программы за 3 года",
             "subtitle": "Срок окупаемости — 14 месяцев"},
            {"type": "quote", "eyebrow": "ПРИНЦИП",
             "quote": "Цифровая трансформация — это не про технологии, а про скорость принятия решений.",
             "attribution": "Из стратегии компании на 2026–2028"},
            {"type": "closing", "eyebrow": "NEO CONSULTING",
             "title": "Спасибо!",
             "subtitle": "Открыты к вопросам и обсуждению следующих шагов",
             "presenter": "Сергей Кузюков", "date": "14 августа 2026"},
        ],
    }


def main() -> int:
    ap = argparse.ArgumentParser(description="Build .pptx from presentation-maker JSON spec")
    ap.add_argument("spec", nargs="?", help="JSON spec file")
    ap.add_argument("out", nargs="?", help="output .pptx path")
    ap.add_argument("--demo", help="write a demo deck to this path and exit")
    args = ap.parse_args()

    if args.demo:
        build(_demo_spec(), Path(args.demo))
        print(f"Demo deck written: {args.demo}")
        return 0

    if not args.spec or not args.out:
        ap.print_help()
        return 2

    try:
        spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    except Exception as e:
        print(f"Failed to read spec {args.spec}: {e}", file=sys.stderr)
        return 1

    out = build(spec, Path(args.out))
    print(f"Deck written: {out} ({len(spec.get('slides', []))} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())